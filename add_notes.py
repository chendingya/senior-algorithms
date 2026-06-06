#!/usr/bin/env python3
"""
Add speaker notes to PPTX file from the lecture script.
"""
from pptx import Presentation
from pptx.util import Inches

# Speaker notes for each slide (extracted from the lecture script)
NOTES = [
    # Slide 1: Cover
    """大家好，我们小组汇报的主题是「游戏中的碰撞检测算法——从四叉树到GJK」。碰撞检测是游戏引擎和物理仿真中最基础也最核心的问题之一。今天我们会沿着「先粗筛、再精算」的工程思路，依次讲解四叉树、Sweep and Prune以及GJK三个算法。""",
    
    # Slide 2: Contents
    """我们的汇报将按照这个结构进行：首先介绍碰撞检测问题的背景和形式化定义，然后依次讲解三个核心算法，最后讨论工业界的应用和总结。""",
    
    # Slide 3: Section Divider - Introduction
    """首先，让我们来看看碰撞检测问题是如何产生的。""",
    
    # Slide 4: Introduction - Problem Background
    """玩家在游戏中看到的是穿模、卡墙这些体验问题，但开发者背后面对的，是碰撞检测这个算法问题。假设场景里有一万个对象，暴力两两检测的复杂度是O(N²)，每帧需要执行约五千万次比较——这在实时渲染中是完全不可接受的。所以真正困难的不是「两个物体能不能碰」，而是「在一万个对象中，哪些对值得做精确判断」。这就是碰撞检测流水线要解决的核心问题。""",
    
    # Slide 5: Formal Definition
    """我们先对问题做形式化定义。碰撞检测的输入是场景中N个物体，每个物体有几何形状和位置；输出是所有相交的物体对。工程上不会直接对所有对做精确检测，而是分成两个阶段。第一阶段叫「粗筛」——用轻量级的近似方法快速排除不可能碰撞的物体对，输出候选集合；第二阶段叫「精算」——只对候选对做精确几何判定。四叉树和SAP属于粗筛阶段，GJK属于精算阶段。这种分层设计是现代游戏引擎的通用架构。""",
    
    # Slide 6: Section Divider - Quadtree
    """接下来我们讲解第一个算法：四叉树。""",
    
    # Slide 7: Quadtree Design Philosophy
    """四叉树的设计思想是「分治」——把二维地图递归划分为四个象限，区域内对象超阈值时继续分裂。直觉是：稀疏区域的对象不必和其他区域比较。查询时只在目标叶节点及相邻区域搜索，大幅缩小范围。这类似给地图分区管理——每个管理员只负责自己地盘。四叉树非常适合二维稀疏场景。""",
    
    # Slide 8: Quadtree Pseudocode & Complexity
    """插入时，对象数未超限就直接存入，超限则分裂后递归插入。查询时，节点与范围不相交则直接剪枝，否则搜索当前对象和子节点。平均情况下，单次查询O(log n)，全场景检测O(n log n)。最坏情况——所有对象集中在同一区域——树会退化，退回到O(n²)。所以四叉树的效果高度依赖空间分布是否均匀。""",
    
    # Slide 9: Quadtree Correctness & Trade-offs
    """四叉树的正确性基于一个不变量：每个对象一定被存储在某个节点中，且四个子节点的空间覆盖了父节点全部区域，不会遗漏。查询时，只要对象在范围内，所在叶节点一定与查询范围相交，保证能被找到。但四叉树有两个核心问题：高动态场景中对象频繁移动，需要不断删除和重新插入，维护代价很高；对象分布极不均匀时，树深急剧增加，丧失加速效果。因此四叉树更适合偏静态的二维场景。""",
    
    # Slide 10: Section Divider - SAP
    """接下来我们讲解第二个算法：Sweep and Prune。""",
    
    # Slide 11: SAP Design Philosophy
    """SAP的设计思想是「降维」——把二维碰撞简化为一维区间重叠问题。将物体包围盒投影到坐标轴，排序后扫描找重叠。这个算法最早出自Baraff 1992年博士论文，后由Cohen等人1995年I-COLLIDE系统化。关键洞察有两个：一是利用时间局部性——帧间变化小，插入排序接近O(n)；二是利用不等式传递性——A与B不重叠、B与C不重叠，则A与C一定不重叠，可直接跳过，实现大规模剪枝。""",
    
    # Slide 12: SAP Pseudocode & Complexity
    """SAP分两步：先排序，再扫描。首次帧用快速排序O(n log n)，后续帧序列几乎有序，改用插入排序接近O(n)。扫描时从左到右遍历，利用不等式传递性——若A与B不重叠、B与C不重叠，则A与C一定也不重叠——实现提前break终止内层循环，复杂度O(n + m)，m为实际重叠数。这和暴力法本质不同：SAP复杂度与输出规模挂钩（output-sensitive），碰撞对越少越接近O(n)。最坏情况下对象剧烈穿插完全乱序，才退化为O(n²)。""",
    
    # Slide 13: SAP Correctness & Applicability
    """SAP的正确性基于分离轴定理：两物体空间相交则任意轴投影必有重叠；逆否命题——轴上不重叠则一定不相交——保证筛选安全。SAP的局限：对象密集大量重叠时候选对接近N²，扫描失效；对象尺寸差异大时大区间吞没小区间，同样退化。SAP最适合「大部分对象相安无事」的典型游戏场景。""",
    
    # Slide 14: Section Divider - GJK
    """接下来我们讲解第三个算法：GJK算法。""",
    
    # Slide 15: GJK Design Philosophy
    """GJK的核心数学基础是Minkowski差：把A中每个点减去B中每个点，形成一个新的点集。关键性质：如果A和B有公共点p，则p-p=0，原点一定在Minkowski差中。反过来，若原点在Minkowski差中，说明存在a属于A、b属于B使得a=b，A和B相交。这样「两物体是否相交」的问题，就转化为「Minkowski差是否包含原点」。GJK不显式构造完整的Minkowski差——那太昂贵——而是通过迭代逼近在Minkowski差的边界上找关键点。""",
    
    # Slide 16: GJK Algorithm Flow
    """GJK的每一步只做一件事：沿当前搜索方向，在Minkowski差上找一个支持点。如果支持点没有越过原点，说明Minkowski差不含原点，直接判定不相交。越过了就加入单纯形，检查是否包围原点。2D场景单纯形依次是点、线段、三角形——三角形包围原点即判定相交。GJK的精妙之处在于不需构造完整Minkowski差，每次只需对两个凸体各做一次极值查询。""",
    
    # Slide 17: GJK Complexity Analysis
    """GJK的复杂度有一个关键事实：2D单纯形最多三角形（3个点），迭代≤3次；3D最多四面体，≤4次。迭代次数是常数，真正代价在Support函数。朴素实现需遍历所有顶点O(n)，预处理可降到O(log n)。更关键的是增量模式：物理引擎用上一帧单纯形作初始猜测，帧间变化小时通常1-2次迭代就收敛——使GJK在实际游戏中接近O(1)。空间复杂度恒为O(1)。正是这种稳定性和小开销使GJK成为实时碰撞检测的首选。""",
    
    # Slide 18: GJK Correctness
    """GJK的正确性需要三件事。第一，等价转化已证明「相交」等价于「原点在Minkowski差中」。第二，Support函数每次取搜索方向的极值点，结合Johnson距离子算法，保证单纯形到原点的距离单调不增，在有限维空间中必然收敛。第三，判定的完备性：返回true时单纯形确实包围了原点；返回false时新支持点证明整个Minkowski差都在原点的某一侧，不可能包含原点。整个过程无需穷举所有顶点。""",
    
    # Slide 19: GJK Example
    """让我们用具体例子走一遍GJK。正方形A中心在原点，B中心在(1.5, 0)，X方向重叠0.5，应判定相交。GJK从初始方向(1,0)取支持点(0.5, 1)，再改方向指向原点取第二个点(-3.5, -2)。两点线段不包围原点，取第三个点形成三角形后，检查原点在三角形内部，判定相交。全过程只调用3次Support函数——每次只需对两个正方形各找一次极值——远比穷举所有顶点对高效。""",
    
    # Slide 20: Section Divider - Industry & Summary
    """最后，让我们看看工业界如何组合使用这些算法。""",
    
    # Slide 21: Industry Combinations
    """工业界不同引擎有不同组合策略。Unity使用PhysX，粗筛支持SAP、MBP和ABP三种模式——MBP适合静态场景，ABP自动选择最优策略。Box2D v3使用动态AABB树，由Erin Catto精心设计。Jolt Physics已被Godot 4内置为默认3D物理引擎。精算阶段几乎都是GJK加EPA的组合。工程中的关键不是押宝单一算法，而是根据场景选择合适组合。""",
    
    # Slide 22: Comprehensive Comparison
    """综合来看，三种算法各有分工。四叉树用分治做空间索引，适合稀疏二维场景；SAP利用帧间连续性做区间扫描，是高动态场景的首选粗筛；GJK用迭代逼近做精确凸体判定，迭代次数有常数上界。它们的共同设计哲学是「先快再准，先筛后算」。工程中这些算法不是互相替代，而是组合使用——粗筛快速排除大量不可能的对，精算只处理少量候选对。""",
    
    # Slide 23: Extensions
    """从GJK出发有多个延伸方向。第一是EPA——GJK只能告诉我们「碰了还是没碰」，但物理引擎需要知道「碰了多深」，EPA在GJK基础上扩展单纯形来计算穿透深度。第二是连续碰撞检测CCD——物体运动很快时离散帧检测可能漏碰，CCD沿运动轨迹做连续扫掠。第三是非凸体处理——GJK只处理凸体，游戏模型多为凹形，需先凸分解再交给GJK。这三个方向说明GJK是起点而非终点。""",
    
    # Slide 24: Quote / Summary
    """总结一下，碰撞检测的核心思想不是比谁算得准，而是比谁排除得快。粗筛快速缩小范围，精算只处理少量候选对，在实时约束下完成高效检测。""",
    
    # Slide 25: Team Division
    """以上是小组分工情况。""",
    
    # Slide 26: Thank You
    """以上就是我们小组的汇报，感谢聆听。请问老师和同学们有什么问题吗？""",
]

def add_notes_to_pptx(input_path, output_path, notes):
    """Add speaker notes to each slide in the PPTX file."""
    prs = Presentation(input_path)
    
    for i, slide in enumerate(prs.slides):
        if i < len(notes):
            # Add notes to the slide
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes[i]
            print(f"Added notes to slide {i+1}")
        else:
            print(f"No notes for slide {i+1}, skipping")
    
    prs.save(output_path)
    print(f"\nSaved to {output_path}")

if __name__ == "__main__":
    import sys
    
    input_file = sys.argv[1] if len(sys.argv) > 1 else "output/presentation.pptx"
    output_file = sys.argv[2] if len(sys.argv) > 2 else "output/presentation-with-notes.pptx"
    
    add_notes_to_pptx(input_file, output_file, NOTES)
